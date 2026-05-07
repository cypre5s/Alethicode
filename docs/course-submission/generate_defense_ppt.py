from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path(__file__).with_name("alethicode-defense.pptx")

NAVY = RGBColor(30, 39, 97)
ICE = RGBColor(202, 220, 252)
TEAL = RGBColor(2, 128, 144)
MINT = RGBColor(2, 195, 154)
WHITE = RGBColor(255, 255, 255)
DARK = RGBColor(34, 40, 49)
MUTED = RGBColor(95, 105, 125)
LIGHT = RGBColor(246, 248, 252)


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, title, subtitle=None, dark=False):
    color = WHITE if dark else NAVY
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.42), Inches(12.1), Inches(0.7))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = color
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.64), Inches(1.08), Inches(11.8), Inches(0.38))
        st = sub.text_frame
        st.clear()
        q = st.paragraphs[0]
        q.text = subtitle
        q.font.name = "Microsoft YaHei"
        q.font.size = Pt(14)
        q.font.color.rgb = ICE if dark else MUTED


def add_footer(slide, idx):
    box = slide.shapes.add_textbox(Inches(11.8), Inches(7.05), Inches(0.9), Inches(0.25))
    p = box.text_frame.paragraphs[0]
    p.text = f"{idx:02d}"
    p.alignment = PP_ALIGN.RIGHT
    p.font.name = "Consolas"
    p.font.size = Pt(10)
    p.font.color.rgb = MUTED


def add_card(slide, x, y, w, h, title, body, accent=TEAL):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(225, 232, 244)
    shape.line.width = Pt(1)
    bar = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(0.08), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.color.rgb = accent

    t = slide.shapes.add_textbox(Inches(x + 0.22), Inches(y + 0.14), Inches(w - 0.35), Inches(0.32))
    p = t.text_frame.paragraphs[0]
    p.text = title
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = DARK

    b = slide.shapes.add_textbox(Inches(x + 0.22), Inches(y + 0.55), Inches(w - 0.35), Inches(h - 0.7))
    tf = b.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, line in enumerate(body):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(12.2)
        p.font.color.rgb = MUTED
        p.space_after = Pt(4)


def add_bullets(slide, lines, x=0.8, y=1.6, w=11.8, h=4.8, size=20, color=DARK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(12)


def title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, NAVY)
    accent = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.18))
    accent.fill.solid()
    accent.fill.fore_color.rgb = MINT
    accent.line.color.rgb = MINT
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.72), Inches(11.8), Inches(1.35))
    p = box.text_frame.paragraphs[0]
    p.text = "Alethicode"
    p.font.name = "Georgia"
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = WHITE
    sub = slide.shapes.add_textbox(Inches(0.85), Inches(3.1), Inches(11.0), Inches(0.85))
    q = sub.text_frame.paragraphs[0]
    q.text = "面向非计算机专业 Python 初学者的 AI 智能在线评测平台"
    q.font.name = "Microsoft YaHei"
    q.font.size = Pt(24)
    q.font.color.rgb = ICE
    tag = slide.shapes.add_textbox(Inches(0.9), Inches(5.9), Inches(8.5), Inches(0.4))
    r = tag.text_frame.paragraphs[0]
    r.text = "在线评测 · AI Tutor · 课件 RAG · 课堂协作"
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(16)
    r.font.color.rgb = WHITE


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_slide(prs)

    slides = [
        ("项目背景与痛点", "传统 OJ 与通用 AI 都无法完整覆盖初学者学习闭环", [
            ("OJ 反馈太冷", ["只给 AC/WA/RE", "不解释为什么错", "缺少下一步建议"]),
            ("AI 脱离课堂", ["不了解教师课件", "不了解学生历史", "容易泛化或泄题"]),
            ("教师负担重", ["备题、批改、答疑成本高", "学情分析难自动化", "个性化辅导难规模化"]),
        ]),
        ("项目目标", "构建真实判题驱动的智能编程教育闭环", [
            ("评测闭环", ["题库、编辑器、提交、判题、提交详情"]),
            ("AI 导学", ["阶段化 Tutor", "错误诊断", "AC 复盘与迁移"]),
            ("教学闭环", ["课件 RAG", "课堂协作", "学情画像和管理端"]),
        ]),
        ("系统总体架构", "Vue 前端 + Spring Boot 主后端 + Python AI 微服务", [
            ("用户层", ["学生 OJ", "教师课堂", "管理端", "课件问答"]),
            ("业务层", ["Spring Boot 3.5.12", "账号、题库、提交、AI、课堂"]),
            ("支撑层", ["PostgreSQL/pgvector", "Redis、Memgraph、NATS、Temporal"]),
        ]),
        ("核心功能一：在线评测", "保留真实 OJ 判题能力，形成可靠学习信号", [
            ("题目与编辑", ["题目列表、题面、样例", "CodeMirror 多语言编辑"]),
            ("提交判题", ["提交代码到后端", "Judge Server 沙箱评测"]),
            ("结果复盘", ["提交历史", "错误详情", "作为 AI 诊断证据"]),
        ]),
        ("核心功能二：AI Tutor", "围绕学生做题过程提供阶段化辅导", [
            ("七阶段", ["READING → IDEATING", "SCAFFOLDING → CODING", "ERROR_FEEDBACK → AC_REVIEW → TRANSFER"]),
            ("上下文", ["题目、提交、错误", "课件、画像、历史对话"]),
            ("会话能力", ["检查点恢复", "压缩、分叉", "WebSocket 实时状态"]),
        ]),
        ("AI 可信回答机制", "减少幻觉、泄题和脱离课堂的回答", [
            ("EvidencePack", ["收集题目、提交、课件、画像证据"]),
            ("Reflection", ["回答后自检教学性、引用、安全边界"]),
            ("安全与成本", ["Prompt Safety Filter", "Token 用量统计", "缓存与熔断"]),
        ]),
        ("核心功能三：课件 RAG 问答", "让 AI 回答回到教师课件和知识点", [
            ("语言包", ["课件导入", "页面切分", "知识点抽取"]),
            ("RAG 服务", ["LightRAG", "pgvector", "Memgraph 图谱"]),
            ("问答体验", ["创建会话", "带引用回答", "页面预览与反馈"]),
        ]),
        ("核心功能四：课堂协作", "从个人练习扩展到真实教学组织", [
            ("班级作业", ["班级、成员", "课堂、作业", "学生提交"]),
            ("教师效率", ["AI 智能出题", "人工审核", "课堂监控"]),
            ("学情分析", ["掌握度", "误区", "风险检测"]),
        ]),
        ("后台管理与运维", "保障平台长期可维护和可观察", [
            ("管理端", ["用户、题目、测试用例", "判题服务器、公告、配置"]),
            ("AI 管理", ["AI 配置", "质量报告", "Trace 和提示变体"]),
            ("运维观测", ["Actuator", "Prometheus/Jaeger", "Sentry/GlitchTip"]),
        ]),
        ("数据库与部署", "按工程系统方式组织数据和运行环境", [
            ("数据层", ["PostgreSQL + pgvector", "Redis Session", "Flyway 迁移"]),
            ("微服务", ["tutor-graph", "alethicode-rag", "Judge Server"]),
            ("部署", ["Docker Compose", "Nginx", "健康检查与环境变量"]),
        ]),
        ("测试与质量保障", "自动化测试覆盖多层核心契约", [
            ("后端", ["单元测试", "控制器契约", "集成和架构测试"]),
            ("前端", ["Jest 单元/契约", "Playwright E2E", "视觉回归"]),
            ("验收", ["登录、提交、AI、RAG、课堂、管理端"]),
        ]),
        ("项目创新点", "真实学习信号驱动个性化编程辅导", [
            ("教学创新", ["真实判题 + AI 导学", "阶段化 Tutor", "课件引用"]),
            ("工程创新", ["AI 可观测性", "EvidencePack + Reflection", "RAG 与 OJ 融合"]),
            ("课堂价值", ["教师降本", "学情画像", "个性化复盘"]),
        ]),
        ("不足与改进", "当前系统已经完整，但仍需继续工程化打磨", [
            ("不足", ["外部服务依赖多", "AI 质量受模型影响", "演示环境复杂"]),
            ("改进", ["健康检查和备用录屏", "评估集与反馈闭环", "压测与容量规划"]),
            ("交付", ["补充截图", "执行最终测试", "按学院模板排版"]),
        ]),
    ]

    for idx, (title, subtitle, cards) in enumerate(slides, start=2):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide, LIGHT)
        add_title(slide, title, subtitle)
        positions = [(0.75, 1.65), (4.75, 1.65), (8.75, 1.65)]
        for (x, y), (ct, body) in zip(positions, cards):
            add_card(slide, x, y, 3.55, 4.55, ct, body, accent=TEAL if x < 8 else MINT)
        add_footer(slide, idx)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, NAVY)
    add_title(slide, "总结", "在线评测 + AI Tutor + 课件 RAG + 课堂协作", dark=True)
    add_bullets(
        slide,
        [
            "Alethicode 构建了面向初学者的智能编程教育闭环。",
            "系统保留真实 OJ 判题能力，并让 AI 基于真实学习证据进行辅导。",
            "项目具备课程展示价值，也具备继续扩展到真实教学场景的工程基础。",
            "谢谢各位老师，欢迎批评指正。",
        ],
        x=1.0,
        y=1.8,
        w=11.2,
        h=4.8,
        size=24,
        color=WHITE,
    )

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
