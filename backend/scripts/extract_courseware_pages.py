#!/usr/bin/env python3
import json
import sys
from pathlib import Path

from pypdf import PdfReader
from pptx import Presentation


def extract_pdf(path: Path):
    reader = PdfReader(str(path))
    pages = []
    for index, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        pages.append({
            "page_no": index,
            "content": text.strip()
        })
    return {
        "extraction_method": "pypdf",
        "pages": pages
    }


def extract_pptx(path: Path):
    presentation = Presentation(str(path))
    pages = []
    for index, slide in enumerate(presentation.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = (shape.text or "").strip()
                if text:
                    texts.append(text)
        pages.append({
            "page_no": index,
            "content": "\n".join(texts).strip()
        })
    return {
        "extraction_method": "python-pptx",
        "pages": pages
    }


def main():
    if len(sys.argv) != 2:
        raise SystemExit("usage: extract_courseware_pages.py <file_path>")

    path = Path(sys.argv[1]).expanduser().resolve()
    if not path.is_file():
        raise SystemExit(f"file not found: {path}")

    suffix = path.suffix.lower()
    if suffix == ".pdf":
        payload = extract_pdf(path)
    elif suffix == ".pptx":
        payload = extract_pptx(path)
    elif suffix in {".ppt", ".doc", ".docx"}:
        raise SystemExit("unsupported lesson file type")
    else:
        raise SystemExit(f"unsupported file extension: {suffix}")

    print(json.dumps(payload, ensure_ascii=False))


if __name__ == "__main__":
    main()
