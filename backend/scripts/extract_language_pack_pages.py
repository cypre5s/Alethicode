#!/usr/bin/env python3
"""
extract_language_pack_pages.py

Extracts page-level content from a normalized language pack document.
Supports: .pdf, .pptx, .docx

Usage:
    python3 extract_language_pack_pages.py <file_path>

Output (JSON to stdout):
    {
        "extraction_method": "pypdf",
        "pages": [
            {"page_no": 1, "content": "..."},
            ...
        ]
    }
"""

import json
import sys
from pathlib import Path

from pypdf import PdfReader


def extract_pdf(path: Path) -> dict:
    reader = PdfReader(str(path))
    pages = []
    for index, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        pages.append({
            "page_no": index,
            "content": text.strip(),
        })
    return {
        "extraction_method": "pypdf",
        "pages": pages,
    }


def extract_pptx(path: Path) -> dict:
    from pptx import Presentation

    presentation = Presentation(str(path))
    pages = []
    for index, slide in enumerate(presentation.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = (shape.text or "").strip()
                if text:
                    texts.append(text)
        pages.append({
            "page_no": index,
            "content": "\n".join(texts).strip(),
        })
    return {
        "extraction_method": "python-pptx",
        "pages": pages,
    }


def extract_docx(path: Path) -> dict:
    from docx import Document

    doc = Document(str(path))
    full_text_parts = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            full_text_parts.append(text)

    full_text = "\n".join(full_text_parts)

    pages = [{
        "page_no": 1,
        "content": full_text,
    }]
    return {
        "extraction_method": "python-docx",
        "pages": pages,
    }


def main():
    if len(sys.argv) != 2:
        raise SystemExit("usage: extract_language_pack_pages.py <file_path>")

    path = Path(sys.argv[1]).expanduser().resolve()
    if not path.is_file():
        raise SystemExit(f"File not found: {path}")

    suffix = path.suffix.lower()
    if suffix == ".pdf":
        payload = extract_pdf(path)
    elif suffix == ".pptx":
        payload = extract_pptx(path)
    elif suffix == ".docx":
        payload = extract_docx(path)
    else:
        raise SystemExit(f"Unsupported file extension: {suffix}")

    print(json.dumps(payload, ensure_ascii=False))


if __name__ == "__main__":
    main()
