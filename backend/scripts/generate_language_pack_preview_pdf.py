#!/usr/bin/env python3
"""
generate_language_pack_preview_pdf.py

Builds a simple multi-page preview PDF for language pack documents.
Currently optimized for .pptx courseware so initialization does not depend on
system LibreOffice availability.

Usage:
    python3 generate_language_pack_preview_pdf.py <input_path> <output_pdf_path>
"""

from __future__ import annotations

import sys
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation


PAGE_WIDTH = 1240
PAGE_HEIGHT = 1754
PAGE_MARGIN = 96
HEADER_GAP = 28
BODY_GAP = 16
LINE_SPACING = 10
BACKGROUND_COLOR = "#fffdf9"
TITLE_COLOR = "#0f172a"
META_COLOR = "#64748b"
BODY_COLOR = "#1f2937"
EMPTY_TEXT = "本页未提取到可显示正文。"
FONT_CANDIDATES = [
    "/home/cypress/.fonts/simhei.ttf",
    "/home/cypress/.fonts/wqy-microhei.ttc",
    "/home/cypress/.local/share/fonts/wqy-microhei.ttc",
    "/mnt/c/Windows/Fonts/simhei.ttf",
    "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
]


def extract_pptx_pages(path: Path) -> list[dict[str, str | int]]:
    presentation = Presentation(str(path))
    pages: list[dict[str, str | int]] = []
    for index, slide in enumerate(presentation.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = (shape.text or "").strip()
                if text:
                    texts.append(text)
        pages.append({
            "page_no": index,
            "content": "\n".join(texts).strip(),
        })
    return pages


def load_font(size: int) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    for candidate in FONT_CANDIDATES:
        path = Path(candidate)
        if not path.is_file():
            continue
        try:
            return ImageFont.truetype(str(path), size=size)
        except OSError:
            continue
    return ImageFont.load_default()


def wrap_text(draw: ImageDraw.ImageDraw, text: str, font: ImageFont.ImageFont, max_width: int) -> list[str]:
    normalized = text.replace("\r\n", "\n").replace("\r", "\n").strip()
    if not normalized:
        return [EMPTY_TEXT]

    lines: list[str] = []
    for paragraph in normalized.split("\n"):
        content = paragraph.strip()
        if not content:
            lines.append("")
            continue

        current = ""
        for char in content:
            tentative = current + char
            width = draw.textbbox((0, 0), tentative, font=font)[2]
            if current and width > max_width:
                lines.append(current)
                current = char
            else:
                current = tentative
        if current:
            lines.append(current)
    return lines or [EMPTY_TEXT]


def render_page(document_title: str, page_no: int, content: str,
                title_font: ImageFont.ImageFont, meta_font: ImageFont.ImageFont,
                body_font: ImageFont.ImageFont) -> Image.Image:
    image = Image.new("RGB", (PAGE_WIDTH, PAGE_HEIGHT), BACKGROUND_COLOR)
    draw = ImageDraw.Draw(image)

    y = PAGE_MARGIN
    title = document_title[:80]
    draw.text((PAGE_MARGIN, y), title, fill=TITLE_COLOR, font=title_font)
    y += title_font.size + HEADER_GAP

    meta = f"第 {page_no} 页"
    draw.text((PAGE_MARGIN, y), meta, fill=META_COLOR, font=meta_font)
    y += meta_font.size + BODY_GAP

    max_width = PAGE_WIDTH - PAGE_MARGIN * 2
    lines = wrap_text(draw, content, body_font, max_width)
    line_height = body_font.size + LINE_SPACING
    max_body_bottom = PAGE_HEIGHT - PAGE_MARGIN

    for index, line in enumerate(lines):
        if y + line_height > max_body_bottom:
            if index == 0:
                draw.text((PAGE_MARGIN, y), EMPTY_TEXT, fill=BODY_COLOR, font=body_font)
            else:
                draw.text((PAGE_MARGIN, y), "……", fill=BODY_COLOR, font=body_font)
            break
        draw.text((PAGE_MARGIN, y), line, fill=BODY_COLOR, font=body_font)
        y += line_height

    return image


def build_preview_pdf(input_path: Path, output_pdf_path: Path) -> None:
    suffix = input_path.suffix.lower()
    if suffix != ".pptx":
        raise SystemExit(f"Unsupported preview generation source: {suffix}")

    pages = extract_pptx_pages(input_path)
    if not pages:
        raise SystemExit("No pages extracted from PPTX")

    title_font = load_font(42)
    meta_font = load_font(26)
    body_font = load_font(28)

    document_title = input_path.stem
    images = [
        render_page(document_title, int(page["page_no"]), str(page["content"]), title_font, meta_font, body_font)
        for page in pages
    ]

    output_pdf_path.parent.mkdir(parents=True, exist_ok=True)
    first_page = images[0]
    remaining_pages = images[1:]
    first_page.save(
        output_pdf_path,
        "PDF",
        resolution=150.0,
        save_all=True,
        append_images=remaining_pages,
    )


def main() -> None:
    if len(sys.argv) != 3:
        raise SystemExit("usage: generate_language_pack_preview_pdf.py <input_path> <output_pdf_path>")

    input_path = Path(sys.argv[1]).expanduser().resolve()
    output_pdf_path = Path(sys.argv[2]).expanduser().resolve()

    if not input_path.is_file():
        raise SystemExit(f"Input file not found: {input_path}")

    build_preview_pdf(input_path, output_pdf_path)


if __name__ == "__main__":
    main()
